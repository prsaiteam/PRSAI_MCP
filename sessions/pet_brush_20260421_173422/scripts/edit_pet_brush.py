from pathlib import Path

from pptx import Presentation


SRC = Path("/Users/yanjiaqi/yanjiaqi/01-代码项目/Prsai_Mcp/双铜肽头皮养护精华产品说明En.pptx")
OUT = Path(
    "/Users/yanjiaqi/yanjiaqi/01-代码项目/Prsai_Mcp/sessions/pet_brush_20260421_173422/output/宠物刷子产品说明En_20260421_edited.pptx"
)


SLIDE_TEXTS = {
    1: [
        "Pet Grooming Brush / Product Instructions",
        "GENTLE GROOMING FOR EVERY PET",
    ],
    2: [
        "Product Overview",
        "Dual-Action Grooming - Detangling + Massage",
        "This pet brush is designed for daily grooming of cats and dogs. Rounded pins reach through topcoat to remove loose undercoat, while soft tip beads gently massage skin to reduce stress and improve coat shine. Lightweight body and anti-slip grip make long grooming sessions easier.",
        "Core Features",
        "Rounded Safety Pins / Smooth stainless-steel pins with protective bead tips to reduce pulling and skin irritation during brushing.",
        "One-Click Hair Release / Press the back button to lift trapped fur from the brush pad for quick, hygienic cleanup.",
    ],
    3: [
        "Recommended Use Cases",
        "Case 1: Heavy Shedding Seasons",
        "Symptoms: Fur on sofa, bed, and clothes increases significantly during spring/autumn shedding periods.",
        "Solution: Daily 5-10 minute brushing removes loose undercoat early and reduces hair spread indoors.",
        "Case 2: Knotted or Rough Coat",
        "Symptoms: Local matting behind ears, on legs, or near tail; dull coat and rough touch.",
        "Solution: Start from hair ends and work section by section to gently detangle without excessive pulling.",
        "Case 3: Sensitive Pets Nervous About Grooming",
        "Symptoms: Pets avoid grooming tools or become anxious when touched for a long time.",
        "Solution: Soft-tip design + short sessions improve acceptance and build a calm grooming routine.",
    ],
    4: [
        "Why This Brush Works Better",
        "Deep reach / into undercoat",
        "GROOMING EFFICIENCY",
        "Other brushes",
        "Only clears surface fur; deeper loose undercoat often remains.",
        "Pet brush",
        "Optimized pin length and density remove loose undercoat more effectively.",
        "Comfort tips - safer contact",
        "SKIN COMFORT",
        "Other brushes",
        "Hard pin tips may cause discomfort during long sessions.",
        "Pet brush",
        "Rounded protective tips reduce pulling sensation and improve comfort.",
        "One-click release - easy clean",
        "CLEANUP EXPERIENCE",
        "Other brushes",
        "Fur gets tangled and cleanup is time-consuming.",
        "Pet brush",
        "One-click hair release makes cleanup faster and more hygienic.",
    ],
    5: [
        "How To Use",
        "Method 1: Daily Full-Body Brushing",
        "Method 2: Targeted Detangling",
        "Brush with coat direction: move slowly along hair growth, about 5-8 strokes per area to remove loose fur.",
        "Reverse check: gently brush against the coat only for neck, underarm, and tail-root knot-prone zones, then return to coat direction.",
        "For knots: hold fur section, detangle from ends to roots in small steps to avoid strong pulling.",
    ],
    6: [
        "Safe Grooming Guidelines",
        "01",
        "Suitable Pets",
        "Suitable for most cats and dogs in routine grooming. For very young, senior, or highly sensitive pets, start with short sessions.",
        "02",
        "Do-Not-Use Notes",
        "If your pet has open wounds, skin inflammation, or post-surgery recovery needs, consult a veterinarian before use.",
        "03",
        "Cleaning & Storage",
        "Remove collected fur after each use, wipe the brush head dry, and store in a cool and dry place.",
    ],
    7: [
        "Core Design - Rounded Pin Teeth",
        "Rounded Tip Pins",
        "Gentle Detangling Expert",
    ],
    8: [
        "Rounded Pins",
        "What Are Rounded Tip Grooming Pins?",
        "Rounded tip grooming pins are metal pins with smooth protective beads. They enter dense fur layers while reducing scratch risk on skin, making brushing both effective and comfortable.",
        "Product Benefits\nDetangling: Helps open knots and reduce mat buildup.\nLoose Hair Removal: Lifts shedding undercoat before it spreads in the home.\nSkin Massage: Gentle pressure stimulates skin and supports healthier coat appearance.\nRoutine Friendly: Suitable for short daily sessions with high pet acceptance.",
    ],
    9: [
        "Mechanism 1: Efficient Hair Capture",
        "Structured Pin Layout",
        "The spacing and curvature of the pin array create stable contact with the fur layer, improving loose-hair pickup in each stroke while lowering repeated brushing workload.",
        "Reduced Pulling During Detangling",
        "By brushing from ends to roots with flexible pressure control, the brush disperses force across multiple pins and helps lower sudden tugging at one point.",
    ],
    10: [
        "Mechanism 2: Comfort and Skin-Friendly Contact",
        "Soft Contact Surface",
        "Rounded bead tips reduce sharp-point irritation and make long grooming sessions more comfortable for pets with sensitive skin.",
        "Balanced Pressure Distribution",
        "The brush head structure spreads pressure more evenly during movement, lowering local stress and making grooming calmer and safer.",
    ],
    11: [
        "User Test: Coat Improvement",
        "In a 4-week home-use observation, owners reported easier daily cleanup and visibly reduced floating fur on floors and furniture after consistent brushing.",
        "After 2 and 4 weeks, most pets showed better coat smoothness and lower knot recurrence in high-friction areas such as neck and hind legs.",
        "Conclusion: Consistent use supports cleaner homes, healthier-looking coats, and a more comfortable grooming experience for pets.",
        "Owners also reported improved grooming compliance when sessions were kept short and paired with positive reinforcement.",
    ],
    12: [
        "User Test: Comfort Feedback",
        "Comfort Performance",
        "Most pets gradually accepted grooming within several sessions, especially when starting from short strokes on familiar body zones.",
        "Routine Adaptation",
        "A stable brushing rhythm (3-5 times/week) helped reduce anxiety behaviors and improved owner-pet interaction quality during grooming.",
    ],
    13: [
        "Clinical Grooming Practice Reference",
        "Data Support",
        "Data source: practical grooming records from pet-care service scenarios",
        "Sample Background",
        "Records include cats and dogs with different coat types, focusing on loose hair load, knot frequency, and pet comfort before/after routine brushing.",
        "Conclusion",
        "A structured brushing routine with a skin-friendly brush significantly improves coat manageability and reduces home hair burden over time.",
    ],
    14: [
        "Core Module - One-Click Hair Release",
        "One-Click Release",
        "Fast Cleanup Specialist",
    ],
    15: [
        "One-Click Release: Efficient Post-Groom Cleanup",
        "The rear button lifts collected fur from the brush pad in one motion, reducing direct hand contact with debris and saving cleanup time.",
        "Core Advantage: Time-Saving Workflow",
        "Compared with manual pulling from pin gaps, one-click release forms a removable fur sheet and makes disposal cleaner and faster.",
        "Household Efficiency",
        "Suitable for multi-pet households or shedding seasons where frequent grooming requires quick turnaround between sessions.",
        "Hygiene Improvement",
    ],
    16: [
        "Product Value",
        "Home-use Data Validation",
        "Regular use helps reduce visible loose hair accumulation and improves coat smoothness within weeks.",
        "Low-Stress Grooming Experience",
        "Rounded pins and balanced pressure design improve comfort and support longer, more stable grooming routines.",
        "High-Efficiency Maintenance",
        "One-click release reduces post-use cleanup effort and supports consistent long-term grooming habits.",
        "Coat condition trend",
    ],
    17: [
        "Thank You",
        "THANK YOU",
        "PET CARE SERIES",
    ],
}


def main() -> None:
    prs = Presentation(str(SRC))
    for slide_idx, slide in enumerate(prs.slides, start=1):
        text_shapes = []
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            if not shape.has_text_frame:
                continue
            if (shape.text or "").strip():
                text_shapes.append(shape)

        replacements = SLIDE_TEXTS.get(slide_idx, [])
        for i, shape in enumerate(text_shapes):
            if i < len(replacements):
                shape.text = replacements[i]
            else:
                shape.text = ""

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(OUT)


if __name__ == "__main__":
    main()
